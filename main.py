import os
import warnings
import pandas as pd
import win32com.client
import tabula
import subprocess
import pdfkit
import fitz
import comtypes.client
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches

warnings.filterwarnings("ignore", category=UserWarning)
warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", category=RuntimeWarning)
warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", category=ImportWarning)
warnings.filterwarnings("ignore", category=ResourceWarning)


def convert_doc_to_docx(doc_file):
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(doc_file)
    doc.SaveAs(doc_file.replace(".doc", ".docx"), FileFormat=16)
    doc.Close()
    word.Quit()


def convert_doc_to_html(doc_file):
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(doc_file)
    doc.SaveAs(doc_file.replace(".doc", ".html"), FileFormat=8)
    doc.Close()
    word.Quit()


def convert_xls_to_pdf(xls_file):
    html_file = convert_excel_to_html(xls_file)
    convert_html_to_pdf(html_file)
    os.remove(html_file)


def convert_excel_to_html(excel_file):
    df = pd.read_excel(excel_file)
    html_file = excel_file.replace('.xls', '.html')
    df.to_html(html_file, index=False)
    return html_file


def convert_doc_to_pdf(doc_file):
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(doc_file)
    doc.SaveAs(doc_file.replace(".doc", ".pdf"), FileFormat=17)
    doc.Close()
    word.Quit()


def convert_docx_to_doc(docx_file):
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(docx_file)
    doc.SaveAs(docx_file.replace(".docx", ".doc"), FileFormat=0)
    doc.Close()
    word.Quit()


def convert_docx_to_html(docx_file):
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(docx_file)
    doc.SaveAs(docx_file.replace(".docx", ".html"), FileFormat=10)
    doc.Close()
    word.Quit()


def convert_docx_to_pdf(docx_file):
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(docx_file)
    doc.SaveAs(docx_file.replace(".docx", ".pdf"), FileFormat=17)
    doc.Close()
    word.Quit()


def convert_xlsx_to_xls(xlsx_file):
    excel = win32com.client.Dispatch("Excel.Application")
    wb = excel.Workbooks.Open(xlsx_file)
    wb.SaveAs(xlsx_file.replace(".xlsx", ".xls"), FileFormat=56)
    wb.Close()
    excel.Quit()


def convert_xlsx_to_pdf(xlsx_file):
    excel = win32com.client.Dispatch("Excel.Application")
    wb = excel.Workbooks.Open(xlsx_file)
    wb.ExportAsFixedFormat(0, xlsx_file.replace(".xlsx", ".pdf"))
    wb.Close()
    excel.Quit()


def convert_xls_to_xlsx(xls_file):
    excel = win32com.client.Dispatch("Excel.Application")
    wb = excel.Workbooks.Open(xls_file)
    wb.SaveAs(xls_file.replace(".xls", ".xlsx"), FileFormat=51)
    wb.Close()
    excel.Quit()


def convert_rt_to_docx(rt_file):
    try:
        subprocess.run(['pandoc', '-s', rt_file, '-o', rt_file.replace('.rtf', '.docx')], check=True)
    except subprocess.CalledProcessError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_rt_to_html(rt_file):
    try:
        subprocess.run(['pandoc', '-s', rt_file, '-o', rt_file.replace('.rtf', '.html')])
    except subprocess.CalledProcessError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_rt_to_pdf(rt_file):
    try:
        intermediate_html = rt_file.replace('.rtf', '.html')
        subprocess.run(['pandoc', '-s', rt_file, '-o', intermediate_html])
        pdfkit.from_file(intermediate_html, rt_file.replace('.rtf', '.pdf'))
    except subprocess.CalledProcessError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_html_to_docx(html_file):
    try:
        subprocess.run(['pandoc', '-s', html_file, '-o', html_file.replace('.html', '.docx')])
    except subprocess.CalledProcessError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_html_to_pdf(html_file):
    try:
        pdfkit.from_file(html_file, html_file.replace('.html', '.pdf'))
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_pdf_to_docx(pdf_file):
    try:
        cv = Converter(pdf_file)
        cv.convert(pdf_file.replace('.pdf', '.docx'))
        cv.close()
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_pdf_to_xlsx(pdf_file):
    try:
        dfs = tabula.read_pdf(pdf_file, pages='all', multiple_tables=True)

        if not dfs:
            raise ValueError("The PDF does not have a table structure.")

        excel_path = pdf_file.replace('.pdf', '.xlsx')

        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            for i, df in enumerate(dfs):
                df.to_excel(writer, sheet_name=f'Sheet {i + 1}', index=False)
    except tabula.errors.JavaNotFoundError:
        print("Error: Make sure Java is installed and properly configured to use tabula-py.")
    except ValueError as e:
        print(f"Error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_pdf_to_pptx(pdf_file):
    try:
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[5]
        pdf_doc = fitz.open(pdf_file)
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            text = page.get_text()
            slide = presentation.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = text
        presentation.save(pdf_file.replace('.pdf', '.pptx'))
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_pdf_to_html(pdf_file):
    try:
        doc = fitz.open(pdf_file)
        html_content = ["<!DOCTYPE html><html><head>", "<title>PDF to HTML Conversion</title>", "</head><body>"]
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("html")
            html_content.append(text)
        html_content.append("</body></html>")
        html_file = pdf_file.replace('.pdf', '.html')
        with open(html_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
        doc.close()
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_pptx_to_pdf(pptx_file):
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_file)
        pdf_file = pptx_file.replace('.pptx', '.pdf')
        presentation.SaveAs(pdf_file, 32)
        presentation.Close()
        powerpoint.Quit()
        None
        None
        import gc
        gc.collect()
    except Exception as e:
        print(f"Unexpected error: {e}")


def convert_pptx_to_ppt(pptx_file):
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_file)
        ppt_file = pptx_file.replace('.pptx', '.ppt')
        presentation.SaveAs(ppt_file, FileFormat=1)
        presentation.Close()
        powerpoint.Quit()
        del presentation, powerpoint
    except Exception as e:
        print(f"Unexpected error: {e}")


def display_menu():
    print("Select a file type to convert:")
    print("1. DOC")
    print("2. DOCX")
    print("3. XLSX")
    print("4. XLS")
    print("5. RT")
    print("6. HTML")
    print("7. PDF")
    print("8. PPTX")
    print("9. PPT")
    choice = input("Enter your choice (1-9): ")
    return choice


def main():
    choice = display_menu()
    file_types = {1: "doc", 2: "docx", 3: "xlsx", 4: "xls", 5: "rt", 6: "html", 7: "pdf", 8: "pptx", 9: "ppt"}

    file_type = file_types.get(int(choice))
    if file_type is None:
        print("Invalid choice!")
        return

    file_path = input("Enter the path of the file: ")

    if not os.path.exists(file_path):
        print("File not found!")
        return

    # Define conversion options for each file type
    conversions = {
        "doc": {"docx": convert_doc_to_docx, "html": convert_doc_to_html, "pdf": convert_doc_to_pdf},
        "docx": {"doc": convert_docx_to_doc, "html": convert_docx_to_html, "pdf": convert_docx_to_pdf},
        "xlsx": {"xls": convert_xlsx_to_xls, "pdf": convert_xlsx_to_pdf},
        "xls": {"xlsx": convert_xls_to_xlsx, "pdf": convert_xls_to_pdf},
        "rt": {"docx": convert_rt_to_docx, "html": convert_rt_to_html, "pdf": convert_rt_to_pdf},
        "html": {"docx": convert_html_to_docx, "pdf": convert_html_to_pdf},
        "pdf": {"docx": convert_pdf_to_docx, "xlsx": convert_pdf_to_xlsx, "pptx": convert_pdf_to_pptx,
                "html": convert_pdf_to_html},
        "pptx": {"pdf": convert_pptx_to_pdf, "ppt": convert_pptx_to_ppt},
        "ppt": {"pptx": convert_pptx_to_ppt},
    }

    # Retrieve conversion options for the selected file type
    conversion_options = conversions.get(file_type)
    if conversion_options is None:
        print("Conversion options not available for this file type!")
        return

    # Display available conversion options for the selected file type
    print(f"Conversion options for {file_type.upper()}:")
    for target_format in conversion_options.keys():
        print(f"- {target_format.upper()}")

    # Get user's target format choice
    target_format = input("Enter the target format: ").lower()

    # Retrieve conversion function for the target format
    conversion_function = conversion_options.get(target_format)
    if conversion_function is None:
        print("Invalid target format!")
        return

    # Perform conversion
    conversion_function(file_path)
    print("Conversion complete!")


if __name__ == "__main__":
    main()







         
      